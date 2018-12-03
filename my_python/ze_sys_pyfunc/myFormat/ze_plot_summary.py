
# ze plot functions
# create on Fri Sep 21 20:57:01 UTC 2018
# created by Jiaze He 

# revised on Tue Nov 20 16:44:27 UTC 2018
# add one_signal_plot

import numpy as np

import argparse, os, sys
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec

# textwrap for plotting text
import textwrap

# generating PPT slides
from pptx import Presentation
from pptx.util import Inches

from seisflows.tools.graphics import plot_vector, plot_section, _cscale, get_regular_ticks

from seisflows.config import config, loadpy, tilde_expand, Dict
from seisflows.tools import unix

import datetime

import glob
import os

from scipy import ndimage


def textplot(struct_variable,filename=None,flag_close=None,flag_save_fig=None):
    
    if filename is None:
	filename = 0

    if flag_close is None:
	flag_close = 1 

    if flag_save_fig is None:
	flag_save_fig = 1

    attributes = [a for a in dir(struct_variable)
                  if not (a.startswith('__') and a.endswith('__'))
                  and not (a == 'kernel')]

    fig = plt.figure(figsize=(10, 7)) 
    #fig.suptitle('bold figure suptitle', fontsize=14, fontweight='bold')

    ax = fig.add_subplot(111)
    fig.subplots_adjust(top=0.85)
    ax.set_title(str('parameters'))

    ax.set_xlabel('xlabel')
    ax.set_ylabel('ylabel')

    spc = 0.26

    for i in range(len(attributes)):
        offset = spc * (i) 
        str_temp = attributes[i] + ' = ' + str((getattr(struct_variable, attributes[i])))
        #str_temp = attributes[i] + ' = ' + str("%.5f" % (getattr(struct_variable, attributes[i])))
       ###ax.text(0.3, 9.7-offset,textwrap.shorten(str_temp, 110), fontsize=11)
        if (i < 35):
           ax.text(0.3, 9.7-offset, str_temp, fontsize=9)
        else:
           ax.text(5.3, 9.7-offset+35*spc, str_temp, fontsize=9)
 

    ax.axis([0, 10, 0, 10])
    #plt.axis('off')
    if flag_save_fig==1:
        plt.savefig(filename,format='png', dpi=500)
    #plt.show()
    if flag_close == 1:
        plt.close()



def add_slide_ze(img_path,filename_pptx,left_start=None,top_start=None,width=None,height=None):

    if left_start is None:
       left_start = 0

    if top_start is None:
       top_start = 0.9  

    if width is None:
       width = 9

    if height is None:
       height = 5.7

    prs = Presentation(filename_pptx)

    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(left_start)
    top = Inches(top_start)
    width = Inches(width)
    height = Inches(height)


    pic = slide.shapes.add_picture(img_path, left, top,width=width,height=height)

    # left = Inches(5)
    # height = Inches(5.5)
    # pic = slide.shapes.add_picture(img_path, left, top, height=height)

    prs.save(filename_pptx)


def plot_section_ze(data, y_step_star = None,y_step_end = None, ax=None, cmap='seismic', clip=100, title='', x_interval=1.0, y_interval=1.0):
    
    if (y_step_star == None) and (y_step_end == None): 
        print('')
    else:
        data = data[y_step_star:y_step_end,:]    
        
    y_interval = 4
    x_interval = 4
    
    nt = data.shape[0]
    nrec = data.shape[1]

    d_aspect = nrec / float(nt)
    
    fsize = 6 
    scale_factor = 1.5
    
#    offsets = np.arange(0,nrec,1)
    
    
    if ax is None:
        fig, ax = plt.subplots(figsize=(fsize, scale_factor*fsize))

    im = ax.imshow(data, aspect=scale_factor*d_aspect, clim=_cscale(data, clip=clip))
    im.set_cmap(cmap)
    #cbar=plt.colorbar(fraction=0.046, pad=0.04,orientation='vertical')
    # labels
    ax.set_title(title)
    ax.set_xlabel('Receiver number')
    ax.set_ylabel('Time steps')

def one_signal_plot(x1,y1,figtitle = None, x_label = None, y_label=None,width=None,height=None):
    if (figtitle is None):
	figtitle = 'a signal'
    if (x_label is None):
        x_label = str(r'time ($\mu s$)') 
    if (y_label is None):
        y_label = ' '
    if (width is None):
        width = 2 
    if (height is None):
        height = 1.2

    plt.title(figtitle) 
    plt.plot(x1,y1)
    plt.xlabel(x_label)
    plt.ylabel(y_label)
    plt.tight_layout(rect=[0, 0, width, height])
    plt.show()

def txt_signal_load(txtfile=None):
    if (txtfile is None):
	txtfile="OUTPUT_FILES/plot_source_time_function.txt"
    print('load a signal')
    stf_load = np.loadtxt(txtfile)
    return stf_load
    #one_signal_plot(stf_load[:,0],stf_load[:,1])
    



